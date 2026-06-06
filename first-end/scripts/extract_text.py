#!/usr/bin/env python3
"""
全格式文档提取器 —— 文本 + 图片引用 + 批量目录模式
依赖: pip install pdfplumber python-pptx python-docx

功能:
  单文件 — python extract_text.py file.pdf
  批量   — python extract_text.py --dir ./course_materials/
  文本提取 — 100+ 格式
  旧格式转换 — Windows COM (PowerPoint/Word)，零额外依赖
  图片提取 — 可选用 --img-dir 导出图片，AI 看图后翻译为文字

支持格式:
  纯文本 .txt .md .csv .log .py .java .c .js .ts .json .xml .yaml ... (100+)
  文档   .pdf .pptx .ppt .docx .doc .xlsx .xls .html .rtf .epub .odt .ods .odp
  邮件   .eml
"""

import argparse
import json
import os
import re
import subprocess
import sys
import tempfile
from pathlib import Path

# ── Unicode 安全输出 ─────────────────────────────────────────
if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
if hasattr(sys.stderr, "reconfigure"):
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")


def _safe_filename(path: str) -> str:
    return re.sub(r'[<>:"/\\|?*\s]+', '_', Path(path).stem)


def _truncate(text: str, limit: int) -> str:
    if len(text) > limit:
        text = text[:limit] + f"\n\n... [截断，共 {len(text)} 字符]"
    return text


def _status(msg: str):
    """输出进度到 stderr，不污染 stdout"""
    print(f"[first-end] {msg}", file=sys.stderr, flush=True)


def _read_plain(filepath: str, limit: int = 100000, encoding: str = None) -> str:
    # 用户指定编码 → 直接使用
    if encoding:
        with open(filepath, "r", encoding=encoding) as f:
            return _truncate(f.read(), limit)

    # Tier 1: BOM 检测（UTF-8/UTF-16 BOM 100% 准确）
    with open(filepath, "rb") as f:
        raw_start = f.read(4)
    if raw_start[:3] == b'\xef\xbb\xbf':
        with open(filepath, "r", encoding="utf-8-sig") as f:
            return _truncate(f.read(), limit)
    if raw_start[:2] in (b'\xff\xfe', b'\xfe\xff'):
        with open(filepath, "r", encoding="utf-16") as f:
            return _truncate(f.read(), limit)

    # Tier 2: chardet 智能检测（可选依赖）
    try:
        import chardet
        with open(filepath, "rb") as f:
            raw = f.read(min(100000, os.path.getsize(filepath)))
        result = chardet.detect(raw)
        if result["confidence"] > 0.7:
            with open(filepath, "r", encoding=result["encoding"]) as f:
                return _truncate(f.read(), limit)
    except ImportError:
        pass

    # Tier 3: 硬编码回退列表（覆盖中/日/韩/欧）
    for enc in ["utf-8", "gb18030", "gbk", "gb2312", "big5",
                "shift_jis", "euc-kr", "latin-1"]:
        try:
            with open(filepath, "r", encoding=enc) as f:
                text = f.read()
            break
        except (UnicodeDecodeError, UnicodeError):
            continue
    else:
        raise RuntimeError(f"无法解码: {filepath}")
    return _truncate(text, limit)


def _com_convert_ppt(filepath: str) -> str | None:
    """Windows COM (late bind): .ppt → .pptx"""
    out_path = os.path.join(tempfile.gettempdir(), _safe_filename(filepath) + ".pptx")
    ps = f'''
$_ppt = [Activator]::CreateInstance([Type]::GetTypeFromProgID('PowerPoint.Application', $true))
$_p = $null
try {{
    $_p = $_ppt.Presentations.Open("{filepath}", -1, 0, 0)
    $_p.SaveAs("{out_path}", 24)
}} finally {{
    if ($_p) {{ $_p.Close() }}
    $_ppt.Quit()
}}
'''
    try:
        subprocess.run(["powershell", "-NoProfile", "-Command", ps],
                       capture_output=True, text=True, timeout=120)
        if os.path.exists(out_path) and os.path.getsize(out_path) > 0:
            return out_path
    except (subprocess.TimeoutExpired, FileNotFoundError):
        pass
    return None


def _com_convert_doc(filepath: str) -> str | None:
    """Windows COM (late bind): .doc → .docx"""
    out_path = os.path.join(tempfile.gettempdir(), _safe_filename(filepath) + ".docx")
    ps = f'''
$_word = [Activator]::CreateInstance([Type]::GetTypeFromProgID('Word.Application', $true))
$_doc = $null
try {{
    $_doc = $_word.Documents.Open("{filepath}")
    $_doc.SaveAs2("{out_path}", 16)
}} finally {{
    if ($_doc) {{ $_doc.Close() }}
    $_word.Quit()
}}
'''
    try:
        subprocess.run(["powershell", "-NoProfile", "-Command", ps],
                       capture_output=True, text=True, timeout=120)
        if os.path.exists(out_path) and os.path.getsize(out_path) > 0:
            return out_path
    except (subprocess.TimeoutExpired, FileNotFoundError):
        pass
    return None


# ═══════════════════════════════════════════════════════════════
#  文本提取器 —— 每个返回 list[dict]
# ═══════════════════════════════════════════════════════════════

PLAIN_EXTS = {
    ".txt", ".md", ".markdown", ".csv", ".log", ".rst",
    ".py", ".java", ".c", ".cpp", ".cc", ".cxx", ".h", ".hpp",
    ".js", ".ts", ".jsx", ".tsx", ".mjs", ".vue", ".svelte",
    ".css", ".scss", ".sass", ".less",
    ".json", ".yaml", ".yml", ".toml", ".xml", ".ini", ".cfg", ".conf",
    ".sh", ".bash", ".zsh", ".bat", ".cmd", ".ps1",
    ".sql", ".r", ".rb", ".go", ".rs", ".swift", ".kt", ".kts",
    ".scala", ".lua", ".php", ".pl", ".pm", ".dart", ".ex", ".exs",
    ".vim", ".gitignore", ".dockerfile", ".env", ".properties",
    ".gradle", ".groovy", ".cmake", ".tex", ".bib",
}


def extract_plain(filepath: str, max_pages=None, encoding: str = None) -> list[dict]:
    return [{"page": 1, "total": 1, "text": _read_plain(filepath, encoding=encoding)}]


def extract_html(filepath: str, max_pages=None, encoding: str = None) -> list[dict]:
    try:
        from html.parser import HTMLParser
        class P(HTMLParser):
            def __init__(self):
                super().__init__()
                self.t = []
            def handle_data(self, d):
                s = d.strip()
                if s:
                    self.t.append(s)
        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
            html = f.read()
        p = P()
        p.feed(html)
        text = "\n".join(p.t)
    except Exception:
        html = _read_plain(filepath)
        html = re.sub(r"<(script|style)[^>]*>.*?</\1>", "", html, flags=re.DOTALL | re.I)
        html = re.sub(r"<[^>]+>", " ", html)
        text = re.sub(r"\s+", " ", html).strip()
    return [{"page": 1, "total": 1, "text": text}]


def extract_pdf(filepath: str, max_pages=None, encoding: str = None) -> list[dict]:
    pages = []
    for lib in ["pdfplumber", "PyPDF2", "pypdf"]:
        try:
            if lib == "pdfplumber":
                import pdfplumber
                with pdfplumber.open(filepath) as pdf:
                    total = len(pdf.pages)
                    limit = min(total, max_pages) if max_pages else total
                    for i in range(limit):
                        text = (pdf.pages[i].extract_text() or "").strip()
                        pages.append({"page": i + 1, "total": total, "text": text})
            elif lib == "PyPDF2":
                from PyPDF2 import PdfReader
                reader = PdfReader(filepath)
                total = len(reader.pages)
                limit = min(total, max_pages) if max_pages else total
                for i in range(limit):
                    text = (reader.pages[i].extract_text() or "").strip()
                    pages.append({"page": i + 1, "total": total, "text": text})
            else:
                from pypdf import PdfReader
                reader = PdfReader(filepath)
                total = len(reader.pages)
                limit = min(total, max_pages) if max_pages else total
                for i in range(limit):
                    text = (reader.pages[i].extract_text() or "").strip()
                    pages.append({"page": i + 1, "total": total, "text": text})
            return pages
        except ImportError:
            continue
    raise ImportError("pip install pdfplumber")


def extract_pptx(filepath: str, max_pages=None, encoding: str = None) -> list[dict]:
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("pip install python-pptx")
    prs = Presentation(filepath)
    slides = []
    total = len(prs.slides)
    limit = min(total, max_pages) if max_pages else total
    for i, slide in enumerate(prs.slides):
        if i >= limit:
            break
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(
                        c.text.strip() for c in row.cells if c.text.strip())
                    if row_text:
                        texts.append(row_text)
        slides.append({"page": i + 1, "total": total, "text": "\n".join(texts)})
    return slides


def extract_docx(filepath: str, max_pages=None, encoding: str = None) -> list[dict]:
    try:
        import docx
    except ImportError:
        raise ImportError("pip install python-docx")
    doc = docx.Document(filepath)
    parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(
                c.text.strip() for c in row.cells if c.text.strip())
            if row_text:
                parts.append(row_text)
    return [{"page": 1, "total": 1, "text": "\n".join(parts)}]


def extract_xlsx(filepath: str, max_pages=None, encoding: str = None) -> list[dict]:
    try:
        import openpyxl
    except ImportError:
        raise ImportError("pip install openpyxl")
    wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
    sheets = []
    for name in wb.sheetnames:
        ws = wb[name]
        rows = []
        for idx, row in enumerate(ws.iter_rows(values_only=True)):
            row_text = " | ".join(str(c) for c in row if c is not None)
            if row_text.strip():
                rows.append(row_text)
            if idx > 5000:
                rows.append("... [截断]")
                break
        sheets.append({
            "page": wb.sheetnames.index(name) + 1,
            "total": len(wb.sheetnames),
            "text": f"=== Sheet: {name} ===\n" + "\n".join(rows)
        })
    wb.close()
    return sheets


def extract_ppt(filepath: str, max_pages=None, encoding: str = None) -> list[dict]:
    tmp = _com_convert_ppt(filepath)
    if tmp:
        try:
            pages = extract_pptx(tmp, max_pages)
            _extract_images(tmp, ".pptx", pages)
            return pages
        finally:
            try:
                os.remove(tmp)
            except OSError:
                pass
    raise RuntimeError("提取 .ppt 失败。请确保安装了 PowerPoint，或另存为 .pptx。")


def extract_doc(filepath: str, max_pages=None, encoding: str = None) -> list[dict]:
    tmp = _com_convert_doc(filepath)
    if tmp:
        try:
            pages = extract_docx(tmp, max_pages)
            _extract_images(tmp, ".docx", pages)
            return pages
        finally:
            try:
                os.remove(tmp)
            except OSError:
                pass
    raise RuntimeError("提取 .doc 失败。请确保安装了 Word，或另存为 .docx。")


def extract_xls(filepath: str, max_pages=None, encoding: str = None) -> list[dict]:
    try:
        import xlrd
        wb = xlrd.open_workbook(filepath)
        sheets = []
        for idx, name in enumerate(wb.sheet_names()):
            ws = wb.sheet_by_name(name)
            rows = []
            for r in range(min(ws.nrows, 5000)):
                row_text = " | ".join(
                    str(ws.cell_value(r, c)) for c in range(ws.ncols)
                    if ws.cell_value(r, c) != "")
                if row_text.strip():
                    rows.append(row_text)
            sheets.append({
                "page": idx + 1, "total": wb.nsheets,
                "text": f"=== Sheet: {name} ===\n" + "\n".join(rows)
            })
        return sheets
    except ImportError:
        pass
    raise RuntimeError("提取 .xls 失败。pip install xlrd  或另存为 .xlsx。")


# ── 可选格式（库按需安装）────────────────────────────────────

def extract_rtf(filepath: str, max_pages=None, encoding: str = None) -> list[dict]:
    try:
        from striprtf.striprtf import rtf_to_text
    except ImportError:
        raise ImportError("pip install striprtf")
    with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
        text = rtf_to_text(f.read())
    return [{"page": 1, "total": 1, "text": text}]


def extract_epub(filepath: str, max_pages=None, encoding: str = None) -> list[dict]:
    try:
        import ebooklib
        from ebooklib import epub
    except ImportError:
        raise ImportError("pip install ebooklib")
    from html.parser import HTMLParser
    book = epub.read_epub(filepath)
    chapters = []
    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
        class P(HTMLParser):
            def __init__(self):
                super().__init__()
                self.texts = []
            def handle_data(self, data):
                t = data.strip()
                if t:
                    self.texts.append(t)
        p = P()
        p.feed(item.get_content().decode("utf-8", errors="ignore"))
        if p.texts:
            chapters.append("\n".join(p.texts))
    return [{"page": 1, "total": 1, "text": "\n\n".join(chapters)}]


def extract_eml(filepath: str, max_pages=None, encoding: str = None) -> list[dict]:
    import email
    from email import policy
    with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
        msg = email.message_from_string(f.read(), policy=policy.default)
    parts = [f"From: {msg.get('From', '')}", f"To: {msg.get('To', '')}",
             f"Subject: {msg.get('Subject', '')}", f"Date: {msg.get('Date', '')}", ""]
    for part in msg.walk():
        if part.get_content_type() == "text/plain":
            payload = part.get_payload(decode=True)
            if payload:
                for enc in ["utf-8", "gbk", "latin-1"]:
                    try:
                        parts.append(payload.decode(enc))
                        break
                    except (UnicodeDecodeError, UnicodeError):
                        continue
    return [{"page": 1, "total": 1, "text": "\n".join(parts)}]


def extract_odt(filepath: str, max_pages=None, encoding: str = None) -> list[dict]:
    try:
        from odf import text as odftext, teletype
        from odf.opendocument import load as odf_load
    except ImportError:
        raise ImportError("pip install odfpy")
    doc = odf_load(filepath)
    paragraphs = [teletype.extractText(p).strip()
                  for p in doc.getElementsByType(odftext.P) if teletype.extractText(p).strip()]
    return [{"page": 1, "total": 1, "text": "\n".join(paragraphs)}]


def extract_ods(filepath: str, max_pages=None, encoding: str = None) -> list[dict]:
    try:
        from odf.opendocument import load as odf_load
        from odf.table import Table, TableRow, TableCell
        from odf import teletype
    except ImportError:
        raise ImportError("pip install odfpy")
    doc = odf_load(filepath)
    sheets = []
    for table in doc.getElementsByType(Table):
        rows_text = []
        for row in table.getElementsByType(TableRow):
            cells = [teletype.extractText(c).strip()
                     for c in row.getElementsByType(TableCell)]
            row_text = " | ".join(c for c in cells if c)
            if row_text:
                rows_text.append(row_text)
        sheets.append("\n".join(rows_text))
    return [{"page": 1, "total": 1, "text": "\n\n".join(sheets)}]


def extract_odp(filepath: str, max_pages=None, encoding: str = None) -> list[dict]:
    try:
        from odf.opendocument import load as odf_load
        from odf.draw import Page
        from odf.text import P
        from odf import teletype
    except ImportError:
        raise ImportError("pip install odfpy")
    doc = odf_load(filepath)
    slides = []
    total = 0
    for page in doc.getElementsByType(Page):
        texts = [teletype.extractText(p).strip()
                 for p in page.getElementsByType(P) if teletype.extractText(p).strip()]
        slides.append({"page": total + 1, "total": 1, "text": "\n".join(texts)})
        total += 1
    for s in slides:
        s["total"] = total
    return slides


# ═══════════════════════════════════════════════════════════════
#  图片提取（可选）
# ═══════════════════════════════════════════════════════════════

IMG_DIR = None


def _extract_pptx_images(filepath: str, prefix: str) -> dict:
    """从 PPTX 提取图片 → {slide_index: [path, ...]}"""
    if not IMG_DIR:
        return {}
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(filepath)
        imgs = {}
        for i, slide in enumerate(prs.slides):
            paths = []
            for j, shape in enumerate(slide.shapes):
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img = shape.image
                    ext = img.content_type.split("/")[-1]
                    if ext == "jpeg":
                        ext = "jpg"
                    fname = f"{prefix}_s{i+1}_img{j+1}.{ext}"
                    fpath = os.path.join(IMG_DIR, fname)
                    with open(fpath, "wb") as f:
                        f.write(img.blob)
                    paths.append(fpath)
            if paths:
                imgs[i + 1] = paths
        return imgs
    except Exception:
        return {}


def _extract_docx_images(filepath: str, prefix: str) -> list:
    """从 DOCX 提取图片 → [path, ...]"""
    if not IMG_DIR:
        return []
    try:
        import docx
        doc = docx.Document(filepath)
        paths = []
        for i, rel in enumerate(doc.part.rels.values()):
            if "image" in rel.reltype:
                img = rel.target_part
                ext = os.path.splitext(img.partname)[-1].lstrip(".")
                if ext == "jpeg":
                    ext = "jpg"
                fname = f"{prefix}_img{i+1}.{ext}"
                fpath = os.path.join(IMG_DIR, fname)
                with open(fpath, "wb") as f:
                    f.write(img.blob)
                paths.append(fpath)
        return paths
    except Exception:
        return []


def _extract_pdf_images(filepath: str, prefix: str) -> dict:
    """从 PDF 提取图片 → {page: [path, ...]}"""
    if not IMG_DIR:
        return {}
    try:
        import pdfplumber
        imgs = {}
        with pdfplumber.open(filepath) as pdf:
            for i, page in enumerate(pdf.pages):
                paths = []
                if hasattr(page, "images") and page.images:
                    for j, img_data in enumerate(page.images):
                        try:
                            fname = f"{prefix}_p{i+1}_img{j+1}.png"
                            fpath = os.path.join(IMG_DIR, fname)
                            stream = img_data.get("stream")
                            if stream:
                                with open(fpath, "wb") as f:
                                    f.write(stream.get_data() if hasattr(stream, 'get_data') else stream)
                                paths.append(fpath)
                        except Exception:
                            pass
                if paths:
                    imgs[i + 1] = paths
        return imgs
    except Exception:
        return {}


def _inject_image_refs(pages: list[dict], images: dict, page_key: str = "page"):
    """在文本中插入图片引用标记。"""
    for idx, img_paths in images.items():
        for p in pages:
            if p[page_key] == idx:
                refs = "\n".join(f"[图片: {os.path.basename(fp)}]" for fp in img_paths)
                p["text"] = p["text"] + "\n\n" + refs if p["text"] else refs
                break


def _extract_images(filepath: str, ext: str, pages: list[dict]):
    """根据文件类型提取图片并注入引用。"""
    global IMG_DIR
    if not IMG_DIR:
        return
    prefix = re.sub(r'[<>:"/\\|?*\s]+', '_', Path(filepath).stem)
    try:
        if ext == ".pptx":
            imgs = _extract_pptx_images(filepath, prefix)
            _inject_image_refs(pages, imgs)
        elif ext == ".docx":
            imgs_list = _extract_docx_images(filepath, prefix)
            if imgs_list and pages:
                refs = "\n".join(f"[图片: {os.path.basename(fp)}]" for fp in imgs_list)
                pages[0]["text"] = pages[0]["text"] + "\n\n" + refs if pages[0]["text"] else refs
        elif ext == ".pdf":
            imgs = _extract_pdf_images(filepath, prefix)
            _inject_image_refs(pages, imgs)
    except Exception:
        pass


# ═══════════════════════════════════════════════════════════════
#  路由表
# ═══════════════════════════════════════════════════════════════

ROUTES = {
    **{ext: ("TEXT", extract_plain) for ext in PLAIN_EXTS},
    ".html": ("HTML", extract_html), ".htm": ("HTML", extract_html),
    ".pdf":  ("PDF",  extract_pdf),
    ".pptx": ("PPTX", extract_pptx),
    ".docx": ("DOCX", extract_docx),
    ".xlsx": ("XLSX", extract_xlsx),
    ".ppt":  ("PPT",  extract_ppt),
    ".doc":  ("DOC",  extract_doc),
    ".xls":  ("XLS",  extract_xls),
    ".rtf":  ("RTF",  extract_rtf),
    ".epub": ("EPUB", extract_epub),
    ".eml":  ("EML",  extract_eml),
    ".odt":  ("ODT",  extract_odt),
    ".ods":  ("ODS",  extract_ods),
    ".odp":  ("ODP",  extract_odp),
}


def _hint(ext: str) -> str:
    return {
        ".pdf": "pip install pdfplumber", ".pptx": "pip install python-pptx",
        ".docx": "pip install python-docx", ".xlsx": "pip install openpyxl",
        ".xls": "pip install xlrd", ".rtf": "pip install striprtf",
        ".epub": "pip install ebooklib", ".odt": "pip install odfpy",
        ".ods": "pip install odfpy", ".odp": "pip install odfpy",
        ".ppt": "确保已安装 PowerPoint", ".doc": "确保已安装 Word",
    }.get(ext, "")


# ═══════════════════════════════════════════════════════════════
#  批量目录处理
# ═══════════════════════════════════════════════════════════════

SKIP_DIRS = {".git", ".svn", "node_modules", "__pycache__", ".venv", "venv",
             ".claude", ".vscode", ".idea", "extracted", "images"}


def _is_supported(filepath: str) -> bool:
    ext = Path(filepath).suffix.lower()
    return ext in ROUTES


def process_directory(dirpath: str, output_dir: str = None, max_pages: int = None,
                      img_dir: str = None, output_format: str = "files",
                      encoding: str = None):
    """批量提取目录下所有支持的文件。

    Args:
        dirpath: 要扫描的目录
        output_dir: 输出目录（默认：dirpath/extracted/）
        max_pages: 全局页数限制
        img_dir: 图片导出目录（默认：output_dir/images/）
        output_format: "files" → 每个文件一个 .md；"json" → 单个 JSON
        encoding: 强制编码（仅对纯文本文件生效）

    Returns:
        dict: {"succeeded": [...], "failed": [...], "skipped": int}
    """
    if not os.path.isdir(dirpath):
        raise RuntimeError(f"目录不存在: {dirpath}")

    if output_dir is None:
        output_dir = os.path.join(dirpath, "extracted")
    os.makedirs(output_dir, exist_ok=True)

    if img_dir is None:
        img_dir = os.path.join(output_dir, "images")

    # 图片目录设置
    global IMG_DIR
    if img_dir:
        IMG_DIR = img_dir
        os.makedirs(IMG_DIR, exist_ok=True)

    # 收集所有支持的文件
    _status(f"扫描 {dirpath} ...")
    all_files = []
    for root, dirs, files in os.walk(dirpath):
        dirs[:] = [d for d in dirs if d not in SKIP_DIRS and not d.startswith(".")]
        for fname in files:
            fpath = os.path.join(root, fname)
            if _is_supported(fpath):
                all_files.append(fpath)

    if not all_files:
        _status("未找到支持的文件")
        return {"succeeded": [], "failed": [], "skipped": 0}

    # 按类型排序：PPT/PPTX 优先（课程TOC的主要来源），然后 PDF，然后其他
    _priority = {".pptx": 0, ".ppt": 0, ".pdf": 1, ".docx": 2, ".doc": 2}
    all_files.sort(key=lambda fp: (_priority.get(Path(fp).suffix.lower(), 3),
                                    os.path.basename(fp)))

    # 统计
    ext_counts = {}
    for fp in all_files:
        ext = Path(fp).suffix.lower()
        ext_counts[ext] = ext_counts.get(ext, 0) + 1
    ext_summary = ", ".join(f"{c} {e}" for e, c in sorted(ext_counts.items()))
    _status(f"找到 {len(all_files)} 个文件: {ext_summary}")

    # 逐个处理
    succeeded = []
    failed = []
    total = len(all_files)

    for idx, filepath in enumerate(all_files, 1):
        fname = os.path.basename(filepath)
        ext = Path(filepath).suffix.lower()
        kind, extractor = ROUTES.get(ext, ("?", None))

        if extractor is None:
            failed.append({"file": filepath, "error": f"不支持的格式: {ext}"})
            continue

        _status(f"[{idx}/{total}] {fname} ({kind}) ...")

        try:
            pages = extractor(filepath, max_pages, encoding=encoding)
        except Exception as e:
            _status(f"[{idx}/{total}] {fname} ❌ {e}")
            failed.append({"file": filepath, "error": str(e)})
            continue

        # 提取图片（IMG_DIR 已在函数开头设置）
        _extract_images(filepath, ext, pages)

        # 输出
        base = Path(filepath).stem
        out_path = os.path.join(output_dir, f"{base}_extracted.md")
        if output_format == "files":
            # 避免同名 ppt/pptx 覆盖
            if os.path.exists(out_path):
                ext_suffix = Path(filepath).suffix.lstrip(".").lower()
                out_path = os.path.join(output_dir, f"{base}_{ext_suffix}_extracted.md")
            with open(out_path, "w", encoding="utf-8") as f:
                f.write(f"# [{kind}] {fname}\n")
                f.write(f"# 源文件: {filepath}\n\n")
                for p in pages:
                    f.write(f"## ── 第 {p['page']} 页 / 共 {p['total']} 页 ──\n")
                    f.write(p["text"])
                    f.write("\n\n")
            _status(f"[{idx}/{total}] {fname} ✅ → {out_path}")

        succeeded.append({
            "file": filepath, "type": kind, "pages": len(pages),
            "output": out_path if output_format == "files" else None,
            "pages_data": pages if output_format == "json" else None,
        })

    # JSON 模式统一输出
    if output_format == "json":
        out_path = os.path.join(output_dir, "extracted_all.json")
        with open(out_path, "w", encoding="utf-8") as f:
            json.dump({
                "source_dir": dirpath,
                "total_files": total,
                "succeeded_count": len(succeeded),
                "failed_count": len(failed),
                "files": succeeded,
                "errors": failed,
            }, f, ensure_ascii=False, indent=2)
        _status(f"✅ JSON → {out_path}")

    # 摘要
    skipped = total - len(succeeded) - len(failed)
    _status(f"完成: {len(succeeded)} 成功, {len(failed)} 失败, {skipped} 跳过")
    for err in failed:
        _status(f"  失败: {os.path.basename(err['file'])} — {err['error']}")

    return {"succeeded": succeeded, "failed": failed, "skipped": skipped}


# ═══════════════════════════════════════════════════════════════
#  主入口
# ═══════════════════════════════════════════════════════════════

def main():
    global IMG_DIR
    parser = argparse.ArgumentParser(description="全格式文档文本提取器")
    parser.add_argument("file", nargs="?", help="文件路径（单文件模式）")
    parser.add_argument("--dir", type=str, default=None, help="批量提取目录下所有文件")
    parser.add_argument("--output-dir", type=str, default=None,
                        help="批量模式输出目录（默认：<dir>/extracted/）")
    parser.add_argument("--output-format", choices=["files", "json"], default="files",
                        help="批量输出格式：files=每文件一个.md, json=单个JSON")
    parser.add_argument("--max-pages", type=int, default=None, help="最大页数")
    parser.add_argument("--encoding", type=str, default=None, help="强制编码（纯文本文件）")
    parser.add_argument("--json", action="store_true", help="单文件模式 JSON 输出")
    parser.add_argument("--img-dir", type=str, default=None, help="图片导出目录")
    parser.add_argument("--cleanup-img", action="store_true", help="清理 img-dir")
    args = parser.parse_args()

    # 清理模式
    if args.cleanup_img:
        if args.img_dir and os.path.isdir(args.img_dir):
            import shutil
            shutil.rmtree(args.img_dir, ignore_errors=True)
            _status(f"清理图片目录: {args.img_dir}")
        return

    # 批量目录模式
    if args.dir:
        process_directory(
            dirpath=args.dir,
            output_dir=args.output_dir,
            max_pages=args.max_pages,
            img_dir=args.img_dir,
            output_format=args.output_format,
            encoding=args.encoding,
        )
        return

    # 单文件模式（兼容旧用法）
    if not args.file:
        print("ERROR: 请指定文件路径 或 --dir 目录路径", file=sys.stderr)
        sys.exit(1)

    filepath = os.path.abspath(args.file)
    if not os.path.exists(filepath):
        print(f"ERROR: 文件不存在: {filepath}", file=sys.stderr)
        sys.exit(1)

    ext = Path(filepath).suffix.lower()
    if ext not in ROUTES:
        print(f"ERROR: 不支持的格式: {ext}", file=sys.stderr)
        sys.exit(1)

    # 图片目录
    if args.img_dir:
        IMG_DIR = args.img_dir
        os.makedirs(IMG_DIR, exist_ok=True)

    kind, extractor = ROUTES[ext]

    try:
        pages = extractor(filepath, args.max_pages, encoding=args.encoding)
    except ImportError as e:
        print(f"ERROR: 缺少库。{_hint(ext)}\n详情: {e}", file=sys.stderr)
        sys.exit(1)
    except Exception as e:
        print(f"ERROR: {e}", file=sys.stderr)
        sys.exit(1)

    # 提取图片（可选）
    _extract_images(filepath, ext, pages)

    if args.json:
        print(json.dumps({"file": filepath, "type": kind, "pages": pages},
                         ensure_ascii=False, indent=2))
    else:
        print(f"# [{kind}] {os.path.basename(filepath)}\n")
        for p in pages:
            print(f"## ── 第 {p['page']} 页 / 共 {p['total']} 页 ──")
            print(p["text"])
            print()


if __name__ == "__main__":
    main()
